"""
PowerPoint document parser using python-pptx.
"""

from pathlib import Path
from typing import List, Dict, Any

from ragBaseMaker.parsers.base_parser import BaseParser, ParsedDocument


class PPTXParser(BaseParser):
    """Parser for PowerPoint presentations."""
    
    SUPPORTED_EXTENSIONS = ['.pptx', '.ppt']
    
    def parse(self, file_path: str | Path) -> ParsedDocument:
        """
        Parse a PowerPoint document.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            ParsedDocument with extracted text and metadata
        """
        from pptx import Presentation
        
        path = self._validate_file(file_path)
        
        prs = Presentation(path)
        
        sections: List[Dict[str, Any]] = []
        all_text_parts: List[str] = []
        metadata: Dict[str, Any] = {
            'slide_count': len(prs.slides),
        }
        
        # Extract core properties if available
        if prs.core_properties:
            props = prs.core_properties
            metadata.update({
                'title': props.title or '',
                'author': props.author or '',
                'subject': props.subject or '',
                'created': str(props.created) if props.created else '',
                'modified': str(props.modified) if props.modified else '',
            })
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts: List[str] = []
            slide_notes: str = ''
            
            # Extract text from shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                
                # Extract from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = ' | '.join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip():
                            slide_texts.append(row_text)
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    slide_notes = notes_frame.text.strip()
            
            if slide_texts or slide_notes:
                sections.append({
                    'type': 'slide',
                    'slide_number': slide_num,
                    'content': slide_texts,
                    'notes': slide_notes,
                })
                
                all_text_parts.append(f"\n[Slide {slide_num}]\n")
                all_text_parts.extend(slide_texts)
                
                if slide_notes:
                    all_text_parts.append(f"\n[Speaker Notes]\n{slide_notes}")
        
        content = self._clean_text('\n'.join(all_text_parts))
        
        return ParsedDocument(
            content=content,
            source_path=str(path.absolute()),
            file_type='pptx',
            metadata=metadata,
            title=metadata.get('title') or path.stem,
            sections=sections,
        )

